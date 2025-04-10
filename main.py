import os
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import streamlit as st

# code program
def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    content = ""

    if ext == ".pdf":
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                content += page.extract_text() + "\n"

    elif ext == ".docx":
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            content += para.text + "\n"

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"

    elif ext == ".xlsx":
        wb = openpyxl.load_workbook(file_path)
        for sheet in wb.worksheets:
            content += f"# {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                content += ' | '.join([str(cell) if cell is not None else "" for cell in row]) + "\n"

    else:
        content = "Format file tidak didukung."

    return content

# Streamlit app
st.title("File Reader")
st.write("Upload a file to read its content. Supported formats: PDF, DOCX, PPTX, XLSX.")

upload_file = st.file_uploader("Choose a file", type=["pdf", "docx", "pptx", "xlsx"])

if upload_file is not None:
    try:
        # Save the uploaded file to a temporary location
        temp_path = os.path.join("temp", upload_file.name)
        os.makedirs(os.path.dirname(temp_path), exist_ok=True)
        with open(temp_path, "wb") as f:
            f.write(upload_file.getbuffer())
                
        # Read the file content
        content = read_file(temp_path)
        st.subheader("File Content")
        st.text_area("Content", content, height=500)
    except Exception as e:
        
        # hapus file sementara
        os.remove(temp_path)
        
    except FileNotFoundError:
        st.error("File tidak ditemukan!")
    except Exception as e:
        st.error(f"Terjadi error saat membaca file: {e}")    
    
    
