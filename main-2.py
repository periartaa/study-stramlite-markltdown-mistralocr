import os
import requests
from dotenv import load_dotenv
from PIL import Image
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import pandas as pd
import tempfile
import streamlit as st

# Load environment variables
load_dotenv()
MISTRAL_API_KEY = os.getenv('MISTRAL_API_KEY')
MISTRAL_API_URL = 'https://api.mistral.ai/v1/ocr'

def process_with_mistral_ocr(file_path, file_type):
    headers = {
        'Authorization': f'Bearer {MISTRAL_API_KEY}',
        'Accept': 'application/json'
    }
    
    try:
        with open(file_path, 'rb') as file:
            response = requests.post(
                MISTRAL_API_URL,
                headers=headers,
                files={'file': file},
                data={'type': file_type}
            )
            
            if response.status_code == 200:
                return response.json().get('text', '')
            else:
                st.error(f"Error with Mistral OCR: {response.status_code} - {response.text}")
                return None
    except Exception as e:
        st.error(f"Error processing file with Mistral OCR: {e}")
        return None

def extract_text_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        text = '\n'.join([page.extract_text() or '' for page in reader.pages])
        if text.strip():
            return text
    except:
        pass

    try:
        images = convert_from_path(pdf_path)
        extracted_text = []

        for image in images:
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                image_path = temp_file.name
                image.save(image_path, 'JPEG')

            text = process_with_mistral_ocr(image_path, 'image')
            if text:
                extracted_text.append(text)

            os.unlink(image_path)

        return '\n'.join(extracted_text)
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_image(image_path):
    return process_with_mistral_ocr(image_path, 'image')

def extract_text_from_word(docx_path):
    try:
        doc = Document(docx_path)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])

        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_part = rel.target_part
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                    image_path = temp_file.name
                    with open(image_path, 'wb') as f:
                        f.write(image_part.blob)

                    image_text = extract_text_from_image(image_path)
                    if image_text:
                        text += f"\n[IMAGE CONTENT]\n{image_text}\n"

                os.unlink(image_path)

        return text
    except Exception as e:
        print(f"Error extracting text from Word document: {e}")
        return None

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)

                if shape.shape_type == 13:
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_file:
                        image_path = temp_file.name
                        with open(image_path, 'wb') as f:
                            f.write(shape.image.blob)

                        image_text = extract_text_from_image(image_path)
                        if image_text:
                            text.append(f"[IMAGE CONTENT]\n{image_text}")

                    os.unlink(image_path)

        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return None

def extract_text_from_excel(excel_path):
    try:
        df = pd.read_excel(excel_path, sheet_name=None)
        text = []
        for sheet_name, sheet_data in df.items():
            text.append(f"=== Sheet: {sheet_name} ===")
            text.append(sheet_data.to_string())
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
        return None

def process_file(file_path):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return None

    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext in ('.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff'):
        return extract_text_from_image(file_path)
    elif file_ext in ('.docx', '.doc'):
        return extract_text_from_word(file_path)
    elif file_ext in ('.pptx', '.ppt'):
        return extract_text_from_pptx(file_path)
    elif file_ext in ('.xlsx', '.xls', '.csv'):
        return extract_text_from_excel(file_path)
    else:
        print(f"Type File Tidak Didukung: {file_ext}")
        return None
    
# Streamlit app
st.title("File Reader")
st.write("Upload a file to read its content. Supported formats: PDF, DOCX, PPTX, XLSX.")

upload_file = st.file_uploader("Supported formats: PDF, DOCX, PPTX, XLSX.")

if upload_file is not None:
    try:
        # Save the uploaded file to a temporary location
        temp_path = os.path.join("temp", upload_file.name)
        os.makedirs(os.path.dirname(temp_path), exist_ok=True)
        with open(temp_path, "wb") as f:
            f.write(upload_file.getbuffer())
                
        # Read the file content
        content = process_file(temp_path)
        st.subheader("File Content")
        st.text_area("Content", content, height=500)

        # hapus file sementara
        os.remove(temp_path)
        
    except FileNotFoundError:
        st.error("File tidak ditemukan!")
    except Exception as e:
        st.error(f"Terjadi error saat membaca file: {e}")    
    
    
